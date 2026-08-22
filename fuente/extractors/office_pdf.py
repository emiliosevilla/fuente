import csv
import json
import logging
import re
import time
from pathlib import Path
from typing import Any, Protocol, Tuple

from fuente.extractors.base import BaseExtractor
from fuente.extractors.base import ExtractionResult
from fuente.extractors.macos_vision import (
    MacOSVisionOCR,
    OCRProcessingError,
    OCRUnavailableError,
)
from fuente.extractors.ocr_image import ImageOCRExtractor
from fuente.extractors.policy import ExtractionPolicy

logger = logging.getLogger(__name__)

class OCRPDFBackend(Protocol):
    def extract_pdf(self, path: Path) -> str: ...


class TextAndOfficeExtractor(BaseExtractor):
    """Extractor local con MarkItDown primero y escalado Docling medido."""

    QUALITY_THRESHOLD = 0.6
    IMAGE_EXTENSIONS = {".png", ".jpeg", ".jpg", ".tiff", ".bmp", ".webp"}
    DOCLING_EXTENSIONS = {".pdf", *IMAGE_EXTENSIONS}
    MARKITDOWN_EXTENSIONS = {
        ".txt", ".md", ".pdf", ".docx", ".doc", ".xlsx", ".xls",
        ".pptx", ".ppt", ".msg", ".html", ".htm", *IMAGE_EXTENSIONS,
    }

    SUPPORTED_EXTENSIONS = {
        ".txt", ".md", ".pdf",
        ".docx", ".doc",
        ".xlsx", ".xls",
        ".pptx", ".ppt",
        ".msg", ".csv",
        ".json", ".html", ".htm", *IMAGE_EXTENSIONS,
    }

    def __init__(self, ocr_backend: OCRPDFBackend | None = None) -> None:
        self.ocr_backend = ocr_backend or MacOSVisionOCR()

    def can_handle(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS

    def extract(self, file_path: Path) -> ExtractionResult | Tuple[str, dict[str, Any]]:
        ext = file_path.suffix.lower()
        metadata = {"original_file": file_path.name, "format": ext}

        # CSV y JSON permanecen nativos para conservar su estructura exacta.
        if ext in {".csv", ".json"}:
            content = self._extract_csv(file_path) if ext == ".csv" else self._extract_json(file_path)
            return ExtractionResult(
                content,
                {**metadata, "extraction_method": "native", "extraction_status": "completed"},
            )

        attempts: list[dict[str, Any]] = []
        degradations: list[str] = []

        if ext in self.MARKITDOWN_EXTENSIONS:
            markitdown_started = time.perf_counter()
            markitdown_res = self._try_markitdown(file_path)
            markitdown_score = self._quality_score(file_path, markitdown_res)
            if markitdown_res:
                if markitdown_score >= self.QUALITY_THRESHOLD:
                    return self._completed(
                        markitdown_res,
                        file_path,
                        metadata,
                        "markitdown",
                        attempts,
                        markitdown_started,
                    )
                attempts.append(self._attempt(
                    "markitdown", "rejected", markitdown_res, markitdown_score,
                    "quality_below_threshold", markitdown_started,
                ))
            else:
                reason = "markitdown_unavailable_or_failed"
                degradations.append(reason)
                attempts.append(self._attempt(
                    "markitdown", "failed", None, 0.0, reason, markitdown_started,
                ))

        native_started = time.perf_counter()
        try:
            native_result = self._extract_native(file_path, metadata)
        except Exception as error:
            logger.error(f"Error extrayendo {file_path.name}: {error}")
            native_result = ExtractionResult(
                None,
                metadata,
                "failed",
                f"native_error: {type(error).__name__}: {error}",
            )

        native_content, native_metadata, native_status, native_reason = self._normalize_result(
            native_result, metadata
        )
        native_score = self._quality_score(file_path, native_content)
        native_outcome = "accepted" if native_status == "completed" and native_score >= self.QUALITY_THRESHOLD else "rejected"
        attempts.append(self._attempt(
            "ocr" if ext in self.IMAGE_EXTENSIONS or ext == ".pdf" and native_metadata.get("extraction_method") not in {None, "pdf_text"} else "native",
            native_outcome,
            native_content,
            native_score,
            native_reason or (None if native_outcome == "accepted" else "quality_below_threshold"),
            native_started,
        ))
        if native_outcome == "accepted":
            return self._with_attempt_metadata(
                ExtractionResult(native_content, native_metadata, native_status, native_reason),
                attempts,
                degradations,
            )

        if ext in self.DOCLING_EXTENSIONS:
            docling_started = time.perf_counter()
            docling_res = self._try_docling(file_path)
            docling_score = self._quality_score(file_path, docling_res)
            if docling_res:
                docling_outcome = "accepted" if docling_score >= self.QUALITY_THRESHOLD else "rejected"
                attempts.append(self._attempt(
                    "docling", docling_outcome, docling_res, docling_score,
                    None if docling_outcome == "accepted" else "quality_below_threshold",
                    docling_started,
                ))
                if docling_outcome == "accepted":
                    return self._with_attempt_metadata(
                        ExtractionResult(
                            docling_res,
                            {
                                **native_metadata,
                                "extraction_method": "docling",
                                "extraction_status": "completed",
                            },
                            "completed",
                        ),
                        attempts,
                        degradations,
                        escalation="docling",
                    )
            else:
                reason = "docling_unavailable_or_failed"
                degradations.append(reason)
                attempts.append(self._attempt(
                    "docling", "failed", None, 0.0, reason, docling_started,
                ))

        reason = native_reason or "extraction_quality: no accepted extraction"
        return self._with_attempt_metadata(
            ExtractionResult(native_content, native_metadata, "failed", reason),
            attempts,
            degradations,
        )

    def _extract_native(self, path: Path, metadata: dict[str, Any]) -> ExtractionResult | str:
        ext = path.suffix.lower()
        if ext in {".txt", ".md"}:
            return self._extract_txt(path)
        if ext == ".pdf":
            return self._extract_pdf(path, metadata)
        if ext in self.IMAGE_EXTENSIONS:
            return ImageOCRExtractor(ocr_backend=self.ocr_backend).extract(path)
        if ext in {".docx", ".doc"}:
            return self._extract_docx(path)
        if ext in {".xlsx", ".xls"}:
            return self._extract_xlsx(path)
        if ext in {".pptx", ".ppt"}:
            return self._extract_pptx(path)
        if ext == ".msg":
            return self._extract_msg(path)
        if ext in {".html", ".htm"}:
            return self._extract_html(path)
        return self._extract_fallback(path)

    @staticmethod
    def _quality_score(path: Path, content: str | None) -> float:
        return ExtractionPolicy._score(path, content)[0]

    @staticmethod
    def _attempt(
        engine: str,
        outcome: str,
        content: str | None,
        quality_score: float,
        reason: str | None,
        started_at: float,
    ) -> dict[str, Any]:
        return {
            "engine": engine,
            "outcome": outcome,
            "quality_score": quality_score,
            "reason": reason,
            "duration_ms": round((time.perf_counter() - started_at) * 1000),
            "has_content": bool(content and content.strip()),
        }

    def _completed(
        self,
        content: str,
        path: Path,
        metadata: dict[str, Any],
        engine: str,
        attempts: list[dict[str, Any]],
        started_at: float,
    ) -> ExtractionResult:
        attempts.append(self._attempt(engine, "accepted", content, self._quality_score(path, content), None, started_at))
        return self._with_attempt_metadata(
            ExtractionResult(content, {**metadata, "extraction_method": engine, "extraction_status": "completed"}),
            attempts,
            [],
        )

    @staticmethod
    def _normalize_result(
        result: ExtractionResult | str,
        base_metadata: dict[str, Any],
    ) -> tuple[str | None, dict[str, Any], str, str | None]:
        if isinstance(result, ExtractionResult):
            return result.content, {**base_metadata, **result.metadata}, result.status, result.reason
        return result, {
            **base_metadata,
            "extraction_method": "native",
            "extraction_status": "completed",
        }, "completed", None

    @staticmethod
    def _with_attempt_metadata(
        result: ExtractionResult,
        attempts: list[dict[str, Any]],
        degradations: list[str],
        *,
        escalation: str | None = None,
    ) -> ExtractionResult:
        metadata = {
            **result.metadata,
            "extraction_attempts": attempts,
            "extraction_degradations": degradations,
        }
        if escalation:
            metadata["extraction_escalation"] = escalation
            metadata["extraction_escalation_reason"] = "quality_below_threshold"
        return ExtractionResult(result.content, metadata, result.status, result.reason)

    def _try_docling(self, path: Path) -> str | None:
        """Escala sólo PDF/imagen a Docling cuando el llamador ya midió baja calidad."""
        if path.suffix.lower() not in self.DOCLING_EXTENSIONS:
            return None
        try:
            from docling.document_converter import DocumentConverter
            converter = DocumentConverter()
            result = converter.convert(str(path))
            if result and hasattr(result, "document"):
                md_content = result.document.export_to_markdown()
                if md_content:
                    logger.info(f"Extracción Docling exitosa para {path.name}")
                    return md_content
        except Exception as e:
            logger.debug(f"Docling no disponible o error para {path.name}: {e}")
        return None

    def _try_markitdown(self, path: Path) -> str | None:
        """Convierte sólo una ruta local, con plugins y servicios cloud deshabilitados."""
        try:
            from markitdown import MarkItDown
            md = MarkItDown(enable_plugins=False)
            res = md.convert_local(path)
            if res and res.text_content:
                logger.info(f"Extracción MarkItDown exitosa para {path.name}")
                return res.text_content
        except Exception as error:
            logger.debug(f"MarkItDown no disponible o error para {path.name}: {error}")
        return None

    def _extract_txt(self, path: Path) -> str:
        """Lee archivos de texto probando codificaciones UTF-8, Latin-1 y Windows-1252."""
        for enc in ["utf-8", "latin-1", "cp1252"]:
            try:
                with open(path, "r", encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    def _extract_pdf(self, path: Path, metadata: dict[str, Any]) -> ExtractionResult:
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(path) as pdf:
                page_count = len(pdf.pages)
                for i, page in enumerate(pdf.pages):
                    t = page.extract_text()
                    if t:
                        text_parts.append(f"<!-- Página {i+1} -->\n" + t)
            extracted = "\n\n".join(text_parts).strip()
            metadata = {**metadata, "page_count": page_count, "extraction_method": "pdf_text"}
            if extracted:
                return ExtractionResult(extracted, {**metadata, "extraction_status": "completed"})
        except ImportError as error:
            return ExtractionResult(
                None,
                {**metadata, "extraction_status": "failed"},
                "failed",
                f"pdf_text_unavailable: {error}",
            )
        except Exception as error:
            return ExtractionResult(
                None,
                {**metadata, "extraction_status": "failed"},
                "failed",
                f"pdf_text_error: {type(error).__name__}: {error}",
            )

        try:
            ocr_text = self.ocr_backend.extract_pdf(path).strip()
            metadata["extraction_method"] = getattr(
                self.ocr_backend,
                "last_method",
                getattr(self.ocr_backend, "method", "macos_vision"),
            )
        except OCRUnavailableError as error:
            reason = f"ocr_unavailable: {error}"
            return ExtractionResult(None, {**metadata, "extraction_status": "failed", "extraction_reason": reason}, "failed", reason)
        except OCRProcessingError as error:
            reason = f"ocr_error: {error}"
            return ExtractionResult(None, {**metadata, "extraction_status": "failed", "extraction_reason": reason}, "failed", reason)
        except Exception as error:
            reason = f"ocr_error: {type(error).__name__}: {error}"
            return ExtractionResult(None, {**metadata, "extraction_status": "failed", "extraction_reason": reason}, "failed", reason)
        if not ocr_text:
            reason = "ocr_empty: Vision no devolvió texto"
            return ExtractionResult(None, {**metadata, "extraction_status": "failed", "extraction_reason": reason}, "failed", reason)
        return ExtractionResult(
            ocr_text,
            {**metadata, "extraction_status": "completed"},
        )

    def _extract_docx(self, path: Path) -> str:
        try:
            import docx
            doc = docx.Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            logger.warning(f"Error extrayendo DOCX {path.name}: {e}. Usando lectura de respaldo.")
            return self._extract_fallback(path)

    def _extract_xlsx(self, path: Path) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            output = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                output.append(f"## Hoja: {sheet}\n")
                for row in ws.iter_rows(values_only=True):
                    row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_str.strip(" |"):
                        output.append(row_str)
            return "\n".join(output)
        except Exception as e:
            logger.warning(f"Error extrayendo XLSX {path.name}: {e}. Usando lectura de respaldo.")
            return self._extract_fallback(path)

    def _extract_pptx(self, path: Path) -> str:
        try:
            import pptx
            prs = pptx.Presentation(path)
            output = []
            for i, slide in enumerate(prs.slides):
                output.append(f"## Diapositiva {i+1}\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        output.append(shape.text)
            return "\n\n".join(output)
        except Exception as e:
            logger.warning(f"Error extrayendo PPTX {path.name}: {e}. Usando lectura de respaldo.")
            return self._extract_fallback(path)

    def _extract_msg(self, path: Path) -> str:
        try:
            import extract_msg
            msg = extract_msg.Message(path)
            output = f"**De:** {msg.sender}\n**Para:** {msg.to}\n**Asunto:** {msg.subject}\n**Fecha:** {msg.date}\n\n{msg.body}"
            return output
        except Exception as e:
            logger.warning(f"Error extrayendo MSG {path.name}: {e}. Usando lectura de respaldo.")
            return self._extract_fallback(path)

    def _extract_csv(self, path: Path) -> str:
        """Convierte archivos CSV a tabla Markdown legible autodetectando el delimitador."""
        lines = []
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            sample = f.read(2048)
            f.seek(0)
            delimiter = ","
            if sample:
                try:
                    dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
                    delimiter = dialect.delimiter
                except Exception:
                    pass
            reader = csv.reader(f, delimiter=delimiter)
            for idx, row in enumerate(reader):
                line = " | ".join(row)
                lines.append(f"| {line} |")
                if idx == 0:
                    lines.append("| " + " | ".join(["---"] * len(row)) + " |")
        return "\n".join(lines)

    def _extract_json(self, path: Path) -> str:
        """Formatea un archivo JSON a bloque de código Markdown estructurado."""
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        formatted = json.dumps(data, indent=2, ensure_ascii=False)
        return f"```json\n{formatted}\n```"

    def _extract_html(self, path: Path) -> str:
        """Limpia etiquetas HTML convirtiéndolas a texto Markdown estructurado."""
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            raw_html = f.read()

        # Reemplazar encabezados HTML por encabezados Markdown
        text = re.sub(r"<h1[^>]*>(.*?)</h1>", r"# \1\n", raw_html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<h2[^>]*>(.*?)</h2>", r"## \1\n", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<h3[^>]*>(.*?)</h3>", r"### \1\n", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<p[^>]*>(.*?)</p>", r"\1\n\n", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<br\s*/?>", "\n", text, flags=re.IGNORECASE)
        
        # Eliminar el resto de etiquetas HTML
        clean_text = re.sub(r"<[^>]+>", "", text)
        return clean_text.strip()

    def _extract_fallback(self, path: Path) -> str:
        """Lectura por defecto segura de caracteres binarios/texto."""
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                cleaned = "".join(c for c in content if c.isprintable() or c in "\n\r\t")
                return cleaned if cleaned.strip() else f"[Documento {path.name} binario/procesado]"
        except Exception as e:
            return f"[No se pudo extraer contenido de {path.name}: {str(e)}]"
