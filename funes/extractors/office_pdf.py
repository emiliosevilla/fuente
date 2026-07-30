import csv
import json
import re
import logging
from pathlib import Path
from typing import Tuple, Dict, Any

from funes.extractors.base import BaseExtractor

logger = logging.getLogger(__name__)


class TextAndOfficeExtractor(BaseExtractor):
    """Extractor completo para TXT, PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, MSG, CSV, JSON, HTML con Docling y MarkItDown."""

    SUPPORTED_EXTENSIONS = {
        ".txt", ".md", ".pdf",
        ".docx", ".doc",
        ".xlsx", ".xls",
        ".pptx", ".ppt",
        ".msg", ".csv",
        ".json", ".html", ".htm"
    }

    def can_handle(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS

    def extract(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        ext = file_path.suffix.lower()
        metadata = {"original_file": file_path.name, "format": ext}

        # 1. Intenta primero Docling (IBM) si está instalado
        docling_res = self._try_docling(file_path)
        if docling_res:
            return docling_res, metadata

        # 2. Intenta MarkItDown (Microsoft) si está instalado
        markitdown_res = self._try_markitdown(file_path)
        if markitdown_res:
            return markitdown_res, metadata

        # 3. Extractores específicos de formato en Python
        try:
            if ext in {".txt", ".md"}:
                return self._extract_txt(file_path), metadata
            elif ext == ".pdf":
                return self._extract_pdf(file_path), metadata
            elif ext in {".docx", ".doc"}:
                return self._extract_docx(file_path), metadata
            elif ext in {".xlsx", ".xls"}:
                return self._extract_xlsx(file_path), metadata
            elif ext in {".pptx", ".ppt"}:
                return self._extract_pptx(file_path), metadata
            elif ext == ".msg":
                return self._extract_msg(file_path), metadata
            elif ext == ".csv":
                return self._extract_csv(file_path), metadata
            elif ext == ".json":
                return self._extract_json(file_path), metadata
            elif ext in {".html", ".htm"}:
                return self._extract_html(file_path), metadata
            else:
                return self._extract_fallback(file_path), metadata
        except Exception as e:
            logger.error(f"Error extrayendo {file_path.name}: {e}")
            return f"[Error de extracción en {file_path.name}: {str(e)}]", metadata

    def _try_docling(self, path: Path) -> str | None:
        """Intenta extraer vía Docling (IBM Research) si se encuentra disponible en el entorno."""
        try:
            from docling.document_converter import DocumentConverter
            converter = DocumentConverter()
            result = converter.convert(str(path))
            if result and hasattr(result, "document"):
                md_content = result.document.export_to_markdown()
                if md_content:
                    logger.info(f"Extracción Docling exitosa para {path.name}")
                    return md_content
        except Exception as e:
            logger.debug(f"Docling no disponible o error para {path.name}: {e}")
        return None

    def _try_markitdown(self, path: Path) -> str | None:
        """Intenta extraer vía MarkItDown (Microsoft) si se encuentra disponible."""
        try:
            from markitdown import MarkItDown
            md = MarkItDown()
            res = md.convert(str(path))
            if res and res.text_content:
                logger.info(f"Extracción MarkItDown exitosa para {path.name}")
                return res.text_content
        except Exception:
            pass
        return None

    def _extract_txt(self, path: Path) -> str:
        """Lee archivos de texto probando codificaciones UTF-8, Latin-1 y Windows-1252."""
        for enc in ["utf-8", "latin-1", "cp1252"]:
            try:
                with open(path, "r", encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    def _extract_pdf(self, path: Path) -> str:
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    t = page.extract_text()
                    if t:
                        text_parts.append(f"<!-- Página {i+1} -->\n" + t)
            return "\n\n".join(text_parts) if text_parts else f"[PDF {path.name} sin texto extraíble o escaneado. Utilizar OCR.]"
        except ImportError:
            return self._extract_fallback(path)

    def _extract_docx(self, path: Path) -> str:
        try:
            import docx
            doc = docx.Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            logger.warning(f"Error extrayendo DOCX {path.name}: {e}. Usando lectura de respaldo.")
            return self._extract_fallback(path)

    def _extract_xlsx(self, path: Path) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            output = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                output.append(f"## Hoja: {sheet}\n")
                for row in ws.iter_rows(values_only=True):
                    row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_str.strip(" |"):
                        output.append(row_str)
            return "\n".join(output)
        except Exception as e:
            logger.warning(f"Error extrayendo XLSX {path.name}: {e}. Usando lectura de respaldo.")
            return self._extract_fallback(path)

    def _extract_pptx(self, path: Path) -> str:
        try:
            import pptx
            prs = pptx.Presentation(path)
            output = []
            for i, slide in enumerate(prs.slides):
                output.append(f"## Diapositiva {i+1}\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        output.append(shape.text)
            return "\n\n".join(output)
        except Exception as e:
            logger.warning(f"Error extrayendo PPTX {path.name}: {e}. Usando lectura de respaldo.")
            return self._extract_fallback(path)

    def _extract_msg(self, path: Path) -> str:
        try:
            import extract_msg
            msg = extract_msg.Message(path)
            output = f"**De:** {msg.sender}\n**Para:** {msg.to}\n**Asunto:** {msg.subject}\n**Fecha:** {msg.date}\n\n{msg.body}"
            return output
        except Exception as e:
            logger.warning(f"Error extrayendo MSG {path.name}: {e}. Usando lectura de respaldo.")
            return self._extract_fallback(path)

    def _extract_csv(self, path: Path) -> str:
        """Convierte archivos CSV a tabla Markdown legible autodetectando el delimitador."""
        lines = []
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            sample = f.read(2048)
            f.seek(0)
            delimiter = ","
            if sample:
                try:
                    dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
                    delimiter = dialect.delimiter
                except Exception:
                    pass
            reader = csv.reader(f, delimiter=delimiter)
            for idx, row in enumerate(reader):
                line = " | ".join(row)
                lines.append(f"| {line} |")
                if idx == 0:
                    lines.append("| " + " | ".join(["---"] * len(row)) + " |")
        return "\n".join(lines)

    def _extract_json(self, path: Path) -> str:
        """Formatea un archivo JSON a bloque de código Markdown estructurado."""
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        formatted = json.dumps(data, indent=2, ensure_ascii=False)
        return f"```json\n{formatted}\n```"

    def _extract_html(self, path: Path) -> str:
        """Limpia etiquetas HTML convirtiéndolas a texto Markdown estructurado."""
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            raw_html = f.read()

        # Reemplazar encabezados HTML por encabezados Markdown
        text = re.sub(r"<h1[^>]*>(.*?)</h1>", r"# \1\n", raw_html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<h2[^>]*>(.*?)</h2>", r"## \1\n", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<h3[^>]*>(.*?)</h3>", r"### \1\n", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<p[^>]*>(.*?)</p>", r"\1\n\n", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<br\s*/?>", "\n", text, flags=re.IGNORECASE)
        
        # Eliminar el resto de etiquetas HTML
        clean_text = re.sub(r"<[^>]+>", "", text)
        return clean_text.strip()

    def _extract_fallback(self, path: Path) -> str:
        """Lectura por defecto segura de caracteres binarios/texto."""
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                cleaned = "".join(c for c in content if c.isprintable() or c in "\n\r\t")
                return cleaned if cleaned.strip() else f"[Documento {path.name} binario/procesado]"
        except Exception as e:
            return f"[No se pudo extraer contenido de {path.name}: {str(e)}]"
